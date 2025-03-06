import os
import streamlit as st
from together import Together
from dotenv import load_dotenv
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from components import funding, valuation, sentiment, investment
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
import re

# Load API Key
load_dotenv()
api_key = os.getenv("TOGETHER_API_KEY")
client = Together(api_key=api_key)

# Define theme colors
THEME_COLORS = {
    'primary': RGBColor(27, 9, 107),  # #1b096b
    'secondary': RGBColor(30, 8, 128),  # #1e0880
    'accent': RGBColor(219, 234, 254),  # #dbeafe
    'text': RGBColor(30, 58, 138),  # #1e3a8a
    'white': RGBColor(255, 255, 255)
}

def generate_pitch(name, area, year_founded, problem, solution, revenue_model, funding_requirements, 
                  sentiment, valuation, funding, investment):
    prompt = f"""
    Create a comprehensive and compelling pitch deck for {name}, a {area} startup founded in {year_founded}.
    Focus on the following core aspects and smartly suggest additional details based on the industry and context:

    1. Executive Summary:
    - Brief overview of {name}
    - Vision and mission
    - Key differentiators and unique value proposition

    2. Problem & Solution:
    Problem: {problem}
    Solution: {solution}
    - Elaborate on the market pain points
    - Explain how the solution addresses these challenges
    - Highlight the technological/innovative aspects

    3. Business & Revenue Model:
    {revenue_model}
    - Suggest pricing strategy
    - Define customer segments
    - Outline scalability potential

    4. Funding Requirements & Use:
    {funding_requirements}
    - Specific allocation of funds
    - Growth milestones
    - Expected timeline

    Market Intelligence:
    - Market Sentiment: {sentiment}
    - AI-Predicted Valuation: ${valuation:,.2f}
    - Projected Funding Rounds: {funding:.1f}
    - Investment Insights: {investment}

    Based on this information, also suggest:
    1. A suitable go-to-market strategy
    2. Potential competitive advantages
    3. Ideal team structure and key roles needed
    4. Market expansion opportunities

    Format the pitch to be compelling, data-driven, and investor-ready, with clear sections and actionable insights.
    """
    try:
        response = client.chat.completions.create(
            model="meta-llama/Llama-3.3-70B-Instruct-Turbo",
            messages=[{"role": "user", "content": prompt}]
        )
        return response.choices[0].message.content.strip()
    except Exception as e:
        return f"❌ Error generating pitch: {e}"

def clean_text_for_pdf(text):
    """Clean text for PDF compatibility."""
    # Replace stars and other special characters
    text = text.replace('*', '•')
    text = text.replace(':', '-')
    # Convert any other problematic characters
    return text.encode('latin-1', 'replace').decode('latin-1')

def create_stylish_pdf(pitch_content, company_name):
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    
    # Add cover page
    pdf.add_page()
    pdf.set_fill_color(27, 9, 107)  # #1b096b
    pdf.rect(0, 0, 210, 297, 'F')  # Full page background
    pdf.set_text_color(255, 255, 255)
    pdf.set_font("Arial", style="B", size=30)
    pdf.cell(0, 120, "", ln=True)  # Spacing
    pdf.cell(0, 20, clean_text_for_pdf(company_name), ln=True, align="C")
    pdf.set_font("Arial", size=20)
    pdf.cell(0, 10, "Pitch Deck", ln=True, align="C")
    
    # Add content pages
    pdf.add_page()
    pdf.set_fill_color(255, 255, 255)
    pdf.rect(0, 0, 210, 297, 'F')
    pdf.set_text_color(30, 58, 138)  # #1e3a8a
    pdf.set_font("Arial", style="B", size=12)
    
    # Split content into sections and style them
    sections = pitch_content.split('\n\n')
    for section in sections:
        if section.strip():
            pdf.set_font("Arial", style="B", size=14)
            # Handle special characters and encoding
            title = clean_text_for_pdf(section.split('\n')[0])
            pdf.cell(0, 10, title, ln=True)
            pdf.set_font("Arial", size=12)
            for line in section.split('\n')[1:]:
                # Handle special characters and encoding
                cleaned_line = clean_text_for_pdf(line.strip())
                if cleaned_line.startswith('•'):
                    pdf.cell(10, 8, '')  # Indent bullets
                pdf.multi_cell(0, 8, cleaned_line)
            pdf.ln(5)
    
    pdf_output = f"{company_name.lower().replace(' ', '_')}_pitch_deck.pdf"
    pdf.output(pdf_output)
    return pdf_output

def add_slide_number(slide, slide_number, total_slides):
    """Add slide number to bottom right."""
    txBox = slide.shapes.add_textbox(
        Inches(13.5), Inches(8.3), Inches(2), Inches(0.5)
    )
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = f"Slide {slide_number} of {total_slides}"
    p.alignment = PP_PARAGRAPH_ALIGNMENT.RIGHT
    p.font.size = Pt(12)
    p.font.color.rgb = THEME_COLORS['text']

def create_stylish_ppt(pitch_content, company_name):
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    
    # Create title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    
    # Set background color
    background = title_slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = THEME_COLORS['primary']
    
    # Add company logo placeholder
    logo_placeholder = title_slide.shapes.add_picture(
        "logo.png", Inches(7), Inches(1), 
        width=Inches(2), height=Inches(2)
    )
    
    # Add title
    title_shape = title_slide.shapes.add_textbox(
        Inches(2), Inches(3.5), Inches(12), Inches(2)
    )
    title_frame = title_shape.text_frame
    title_para = title_frame.add_paragraph()
    title_para.text = company_name
    title_para.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    title_para.font.size = Pt(60)
    title_para.font.color.rgb = THEME_COLORS['white']
    title_para.font.bold = True
    
    # Add subtitle
    subtitle_shape = title_slide.shapes.add_textbox(
        Inches(2), Inches(5.5), Inches(12), Inches(1)
    )
    subtitle_frame = subtitle_shape.text_frame
    subtitle_para = subtitle_frame.add_paragraph()
    subtitle_para.text = "Pitch Deck"
    subtitle_para.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    subtitle_para.font.size = Pt(32)
    subtitle_para.font.color.rgb = THEME_COLORS['accent']
    
    # Create content slides
    sections = pitch_content.split('\n\n')
    total_slides = len(sections) + 1  # Including title slide
    
    for idx, section in enumerate(sections, 1):
        if section.strip():
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            
            # Set slide background
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = THEME_COLORS['white']
            
            # Add title with accent bar
            title_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.5), 
                Inches(15), Inches(1.2)
            )
            title_shape.fill.solid()
            title_shape.fill.fore_color.rgb = THEME_COLORS['primary']
            title_shape.line.color.rgb = THEME_COLORS['primary']
            
            # Add title text
            title_box = slide.shapes.add_textbox(
                Inches(1), Inches(0.7), Inches(14), Inches(0.8)
            )
            title_frame = title_box.text_frame
            title_para = title_frame.add_paragraph()
            title_para.text = section.split('\n')[0]
            title_para.font.size = Pt(44)
            title_para.font.color.rgb = THEME_COLORS['white']
            
            # Add content with better formatting
            content_box = slide.shapes.add_textbox(
                Inches(1), Inches(2), Inches(14), Inches(6)
            )
            content_frame = content_box.text_frame
            
            for line in section.split('\n')[1:]:
                para = content_frame.add_paragraph()
                if line.strip().startswith(('•', '-', '*','?','!')):
                    para.level = 1
                    para.text = line.strip().replace('*', '').replace('-', '').replace('?', '').replace('!', '').strip()
                else:
                    para.text = line.strip()
                para.font.size = Pt(24)
                para.font.color.rgb = THEME_COLORS['text']
            
            # Add slide number
            add_slide_number(slide, idx + 1, total_slides)
    
    ppt_output = f"{company_name.lower().replace(' ', '_')}_pitch_deck.pptx"
    prs.save(ppt_output)
    return ppt_output

def fetch_component_value(component, use_component):
    return component.show() if use_component else None

def show():
    st.markdown("""
    <style>
        @import url('https://fonts.googleapis.com/css2?family=Lora:wght@400;500;600;700&display=swap');

        body {
            font-family: 'Lora', serif;
            font-size: 18px;
        }
    </style>
    """, unsafe_allow_html=True)
    st.title("🚀 Startup Pitch Generator")
    
    with st.form("pitch_form"):
        col1, col2 = st.columns(2)
        
        with col1:
            name = st.text_input("Startup Name")
            area = st.text_input("Industry/Area")
            year_founded = st.number_input("Year Founded", min_value=1900, max_value=2100, value=2024, step=1)
        
        with col2:
            use_sentiment = st.checkbox("Include Market Sentiment", value=True)
            use_valuation = st.checkbox("Include Valuation Prediction", value=True)
            use_funding = st.checkbox("Include Funding Prediction", value=True)

        st.subheader("Core Information")
        problem = st.text_area("Problem Statement", height=100)
        solution = st.text_area("Solution", height=100)
        revenue_model = st.text_area("Revenue Model", height=100)
        funding_requirements = st.text_area("Funding Requirements", height=100)
        
        submit_button = st.form_submit_button("Generate Pitch")
    
    if submit_button:
        with st.spinner("Generating your pitch..."):
            # Fetch component values
            sentiment_value = fetch_component_value(sentiment, use_sentiment) or "Not available"
            valuation_value = fetch_component_value(valuation, use_valuation) or 0.0
            funding_value = fetch_component_value(funding, use_funding) or 0.0
            
            pitch_content = generate_pitch(
                name, area, year_founded, problem, solution, revenue_model,
                funding_requirements, sentiment_value, valuation_value, 
                funding_value, "Based on market analysis"
            )
        
        st.subheader("📜 Generated Pitch Deck Content")
        st.markdown(pitch_content)
        
        # Generate both PDF and PPT
        with st.spinner("Creating presentation files..."):
            pdf_path = create_stylish_pdf(pitch_content, name)
            ppt_path = create_stylish_ppt(pitch_content, name)
            
            col1, col2 = st.columns(2)
            with col1:
                with open(pdf_path, "rb") as pdf_file:
                    st.download_button(
                        "📥 Download PDF Pitch Deck",
                        data=pdf_file,
                        file_name=pdf_path,
                        mime="application/pdf",
                        key="pdf_download"
                    )
            
            with col2:
                with open(ppt_path, "rb") as ppt_file:
                    st.download_button(
                        "📥 Download PowerPoint Pitch Deck",
                        data=ppt_file,
                        file_name=ppt_path,
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        key="ppt_download"
                    )

if __name__ == "__main__":
    show()

